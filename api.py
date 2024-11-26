import os
import uvicorn
from fastapi import FastAPI, UploadFile, File, BackgroundTasks
from main import startLabelingFlow
from Agent_config import AgentFactory
from Database import reg_tool
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import dashscope
from dotenv import load_dotenv
import global_vars
from pydantic import BaseModel

load_dotenv()
MULTIMODAL_API_KEY = os.getenv("MULTIMODAL_API_KEY")

app = FastAPI()
shared_data = {}
# 创建保存文件的目录（如果不存在）
os.makedirs("source_file", exist_ok=True)

# 提取ppt文件的文字信息，将图片通过qwen-vl转化为文字信息并保存至output.txt
def process_PPT(file_name):
    print("process PPT:", file_name)
    folder_path = f"./source_file/{file_name}"
    file_path = os.path.join(folder_path, file_name)
    txt_file_path = os.path.join(folder_path, 'output.txt')
    imgs_dir = os.path.join(folder_path, "imgs")
    # 打开PPT文件
    ppt = Presentation(file_path)

    # 创建Markdown文件
    with open(txt_file_path, 'w', encoding='utf-8') as md_file:
        for slide_number, slide in enumerate(ppt.slides, start=1):
            # 提取表格
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        md_file.write(paragraph.text + '\n')

                if shape.has_table:
                    table = shape.table
                    # 写入表格内容
                    for row_index, row in enumerate(table.rows):
                        row_text = '| ' + ' | '.join(cell.text for cell in row.cells) + ' |\n'
                        md_file.write(row_text)
                        if row_index == 0:
                            header = '| ' + ' | '.join(['---'] * len(table.columns)) + ' |\n'
                            md_file.write(header)
                            # 检查形状是否是图片类型

                # 提取图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 获取图片
                    image = shape.image
                    # 获取图片的二进制数据
                    image_bytes = image.blob
                    # 获取图片的文件名（你可以根据需要修改这个）
                    image_filename = f"image_{slide_number}_{shape._element.get('id')}.png"
                    image_file_path = os.path.join(imgs_dir, image_filename)
                    # 保存图片到输出文件夹
                    with open(image_file_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    print(f"Saved {image_filename}")
                    img_msg = vl_api(image_file_path)
                    md_file.write(f"""\n![{img_msg}]({image_file_path})\n""")

                md_file.write('\n')


# 根据文件类型处理(上传的文件保存到source_file文件夹，每个上传文件一个单独文件夹，文件夹名字为上传文件的名字，里面包含：源文件、图片、output.txt)
def process_file(file):
    file_name = file.filename
    print("开始处理 fileName", file_name)
    folder_path = f"./source_file/{file_name}"
    file_path = os.path.join(folder_path, file_name)
    imgs_dir = os.path.join(folder_path, "imgs")
    if not os.path.exists(imgs_dir):
        os.makedirs(imgs_dir)
    # 如果文件夹不存在，则创建
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
    # 保存文件到本地
    with open(file_path, "wb") as f:
        for chunk in file.file:
            f.write(chunk)
    shared_data[file_name] = {
        "fileName": file_name,
        "status": 0
    }
    print("shared_data", shared_data)
    process_PPT(file_name)
    sub_task_name = f"{global_vars.task_name}_process_file"
    global_vars.queueTask.remove(sub_task_name)


def vl_api(img_path):
    messages = [
        {
            "role": "user",
            "content": [
                {"image": f"""file://{img_path}"""},
                {"text": "You are professional assist with information extraction tools\
                 这张图片是来自于富含专业信息的文档截图，按照以下的步骤进行信息提取\
                 第一步，提取图片中的结构信息，尽可能找到文档的层次结构、段落结构、表格结构等信息。\
                 第二步，提取图片中的文本信息，尽可能找到文档中的关键信息、细节信息等。不要改动原始数据的任何信息\
                 第三步，查验核对图片中信息是否正确，不要遗漏任何重要信息。\
                 最后以美观的格式进行输出。"}
            ]
        }
    ]
    response = dashscope.MultiModalConversation.call(
        # 需要替换成自己的api_key=
        api_key=MULTIMODAL_API_KEY,
        model='qwen-vl-max',
        messages=messages
    )
    return response["output"]["choices"][0]["message"]["content"][0]["text"]


@app.post('/run', tags=["后台运行结果接口"],
          description="后台运行结果接口：根据返回值确认后台是否成功运行",
          response_description="返回 code:1 表示成功，code:0 表示失败")
async def deal(task_name:str, flowTask : BackgroundTasks, file: UploadFile = File(...)):
    try:
        # TODO: 上传文件的处理逻辑
        # 需要对执行过程的错误代码进行处理，ValidError或者其他类型的异常处理
        # 检查任务Name是否在执行
        if task_name in global_vars.queueTask:
            return {"code": 0, "message": "任务正在执行中"}
        else:
            global_vars.queueTask.append(task_name)
            global_vars.task_name = task_name
        # 检查文件是否存在
        if os.path.exists(f"./source_file/{file.filename}") is True:
            # rename file
            file.filename = f"{file.filename}_{str(os.path.getctime(f'./source_file/{file.filename}'))}"
            file_name = file.filename
        else:
            file_name = file.filename
        # 处理文件，提取文字信息
        sub_task_name = f"{task_name}_process_file"
        global_vars.queueTask.append(sub_task_name)
        flowTask.add_task(process_file(file))
        # await process_file(file)
        

        # 读取处理好后的文字信息(保存到了output.txt中)
        folder_path = f"./source_file/{file_name}"
        txt_file_path = os.path.join(folder_path, 'output.txt')
        with open(txt_file_path, "r", encoding='utf-8') as f:
            raw_data = f.read()

        # 实例化和注册代理
        agent = AgentFactory()
        reg_tool(agent)
        # 处理数据并保存结果
        sub_task_name = f"{task_name}_startLabelingFlow"
        global_vars.queueTask.append(sub_task_name)
        flowTask.add_task(startLabelingFlow(agent, raw_data))
        global_vars.queueTask.remove(task_name)
    except FileNotFoundError as e:
        return {"code": 0, "message": str(e)}
    except Exception as e:
        # 记录其他类型的异常，返回错误信息
        print(f"发生错误: {str(e)}")
        return {"code": 0, "message": "处理过程中出现错误"}
    # 如果一切正常，返回成功代码 return pydantic.BaseModel
    return {"code": 1, "message": "Success"}

if __name__ == '__main__':
    uvicorn.run("api:app", port=8000, reload=True)
